"""
exporters.py
Модуль экспорта артефактов для проекта Open Executive.
Поддерживает экспорт в PDF и PowerPoint (PPTX).
"""

from typing import List, Dict, Any
from pathlib import Path


def export_to_pdf(markdown_text: str, output_path: str) -> bool:
    """
    Экспортирует Markdown-текст в PDF файл.
    
    Args:
        markdown_text: Текст в формате Markdown.
        output_path: Путь для сохранения PDF файла.
        
    Returns:
        bool: True если экспорт успешен, False иначе.
        
    Note:
        Это заглушка для демонстрации интерфейса.
        Для реальной реализации рекомендуется использовать:
        - markdown + weasyprint
        - или pandoc через subprocess
        - или reportlab для прямого генерирования PDF
    """
    try:
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        
        # Заглушка: сохраняем исходный MD как текстовый файл
        # В продакшене здесь будет реальная конвертация в PDF
        with open(output, 'w', encoding='utf-8') as f:
            f.write("# Exported from Open Executive\n\n")
            f.write("(PDF conversion placeholder)\n\n")
            f.write(markdown_text)
        
        return True
        
    except Exception as e:
        print(f"Error exporting to PDF: {e}")
        return False


def export_to_pptx(structured_data: List[Dict[str, Any]], output_path: str) -> bool:
    """
    Генерирует PowerPoint презентацию из структурированных данных.
    
    Args:
        structured_data: Список словарей со слайдами.
                        Каждый словарь должен иметь ключи:
                        - 'title' (str): Заголовок слайда
                        - 'bullets' (List[str]): Список пунктов для буллитов
        output_path: Путь для сохранения .pptx файла.
        
    Returns:
        bool: True если экспорт успешен, False иначе.
        
    Example:
        >>> data = [
        ...     {"title": "Введение", "bullets": ["Цель проекта", "Задачи"]},
        ...     {"title": "Финансы", "bullets": ["Бюджет: $1M", "ROI: 25%"]}
        ... ]
        >>> export_to_pptx(data, "presentation.pptx")
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("Error: Требуется установить python-pptx: pip install python-pptx")
        return False
    
    try:
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        
        # Создаем презентацию
        prs = Presentation()
        
        # Устанавливаем широкоформатный режим (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for slide_data in structured_data:
            title = slide_data.get('title', 'Без названия')
            bullets = slide_data.get('bullets', [])
            
            # Добавляем слайд с макетом "Заголовок и содержимое"
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Устанавливаем заголовок
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Получаем текстовый фрейм для контента
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()  # Очищаем дефолтный параграф
            
            # Добавляем буллиты
            if bullets:
                # Первый пункт
                first_p = tf.paragraphs[0]
                first_p.text = bullets[0]
                first_p.font.size = Pt(18)
                first_p.level = 0
                
                # Остальные пункты
                for bullet_text in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet_text
                    p.font.size = Pt(18)
                    p.level = 0
            else:
                # Если буллитов нет, добавляем пустой параграф
                p = tf.paragraphs[0]
                p.text = ""
        
        # Сохраняем презентацию
        prs.save(str(output))
        return True
        
    except Exception as e:
        print(f"Error exporting to PPTX: {e}")
        return False


def export_markdown_to_structured_slides(
    markdown_text: str, 
    max_bullets_per_slide: int = 6
) -> List[Dict[str, Any]]:
    """
    Конвертирует Markdown-текст в структурированный формат для слайдов.
    
    Args:
        markdown_text: Текст в формате Markdown.
        max_bullets_per_slide: Максимальное количество пунктов на слайд.
        
    Returns:
        List[Dict[str, Any]]: Список слайдов в формате для export_to_pptx.
    """
    slides = []
    current_title = None
    current_bullets = []
    
    lines = markdown_text.split('\n')
    
    for line in lines:
        line = line.strip()
        
        # Обработка заголовков (H1, H2, H3)
        if line.startswith('###'):
            # Сохраняем предыдущий слайд если есть
            if current_title and current_bullets:
                slides.append({
                    'title': current_title,
                    'bullets': current_bullets[:max_bullets_per_slide]
                })
                current_bullets = []
            
            current_title = line.replace('###', '').strip()
            
        elif line.startswith('##'):
            # Сохраняем предыдущий слайд
            if current_title and current_bullets:
                slides.append({
                    'title': current_title,
                    'bullets': current_bullets[:max_bullets_per_slide]
                })
                current_bullets = []
            
            current_title = line.replace('##', '').strip()
            
        elif line.startswith('#'):
            # Сохраняем предыдущий слайд
            if current_title and current_bullets:
                slides.append({
                    'title': current_title,
                    'bullets': current_bullets[:max_bullets_per_slide]
                })
                current_bullets = []
            
            current_title = line.replace('#', '').strip()
            
        # Обработка списков
        elif line.startswith('- ') or line.startswith('* ') or line.startswith('• '):
            bullet_text = line[2:].strip()
            if bullet_text:
                current_bullets.append(bullet_text)
                
        # Обработка нумерованных списков
        elif line[0].isdigit() and line[1:3] in ['. ', ') ']:
            bullet_text = line.split('. ', 1)[-1].split(') ', 1)[-1].strip()
            if bullet_text:
                current_bullets.append(bullet_text)
                
        # Обычный текст (если есть заголовок)
        elif line and current_title:
            # Разбиваем длинный текст на предложения
            sentences = [s.strip() for s in line.split('.') if s.strip()]
            for sentence in sentences:
                if not sentence.endswith('.'):
                    sentence += '.'
                current_bullets.append(sentence)
    
    # Добавляем последний слайд
    if current_title:
        slides.append({
            'title': current_title,
            'bullets': current_bullets[:max_bullets_per_slide] if current_bullets else ['']
        })
    
    # Если слайдов нет, создаем хотя бы один
    if not slides:
        slides.append({
            'title': 'Open Executive Export',
            'bullets': [markdown_text[:200] + '...' if len(markdown_text) > 200 else markdown_text]
        })
    
    return slides


# Пример использования
if __name__ == "__main__":
    # Тестовые данные
    test_slides = [
        {
            "title": "Стратегия развития Q1 2025",
            "bullets": [
                "Увеличение доли рынка на 15%",
                "Запуск нового продукта в феврале",
                "Оптимизация операционных расходов",
                "Расширение команды на 20 человек"
            ]
        },
        {
            "title": "Финансовые показатели",
            "bullets": [
                "Выручка: $2.5M (+25% YoY)",
                "EBITDA: $450K (маржа 18%)",
                "CAC: $120 (снижение на 15%)",
                "LTV: $1,800 (рост на 30%)"
            ]
        },
        {
            "title": "План действий",
            "bullets": [
                "Январь: Подготовка инфраструктуры",
                "Февраль: Маркетинговая кампания",
                "Март: Анализ результатов и корректировка"
            ]
        }
    ]
    
    # Тест экспорта в PPTX
    output_file = "test_presentation.pptx"
    success = export_to_pptx(test_slides, output_file)
    
    if success:
        print(f"Презентация успешно создана: {output_file}")
        print(f"Количество слайдов: {len(test_slides)}")
    else:
        print("Не удалось создать презентацию")
    
    # Тест конвертации Markdown
    test_md = """
# Обзор проекта

## Цели
- Достичь PMF к Q2
- Привлечь $1M инвестиций

## Риски
- Конкуренция растет
- Регуляторные изменения
"""
    
    structured = export_markdown_to_structured_slides(test_md)
    print(f"\nКонвертировано слайдов из MD: {len(structured)}")
    for slide in structured:
        print(f"  - {slide['title']}: {len(slide['bullets'])} пунктов")
